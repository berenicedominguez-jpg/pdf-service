import os, base64
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import subprocess

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
LOGO_TC  = os.path.join(BASE_DIR, 'assets', 'logo_tc_transparent.png')

ROJO_OSCURO = RGBColor(0xB0,0x10,0x00)
NARANJA     = RGBColor(0xC0,0x60,0x00)
BLANCO      = RGBColor(0xFF,0xFF,0xFF)
AZUL_TITULO = RGBColor(0x8B,0xC4,0xE0)
GRIS_TEXTO  = RGBColor(0x33,0x33,0x33)

W = Inches(13.33)
H = Inches(7.5)

def fondo_rojo(slide):
    bg = slide.shapes.add_shape(1,0,0,W,H)
    bg.fill.solid(); bg.fill.fore_color.rgb = ROJO_OSCURO
    bg.line.fill.background()
    franja = slide.shapes.add_shape(1,0,Inches(6.2),W,Inches(1.3))
    franja.fill.solid(); franja.fill.fore_color.rgb = NARANJA
    franja.line.fill.background()

def fondo_blanco(slide):
    bg = slide.shapes.add_shape(1,0,0,W,H)
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF)
    bg.line.fill.background()
    franja = slide.shapes.add_shape(1,0,Inches(6.8),W,Inches(0.7))
    franja.fill.solid(); franja.fill.fore_color.rgb = NARANJA
    franja.line.fill.background()

def add_confidencial(slide):
    txb = slide.shapes.add_textbox(Inches(4.5),Inches(6.85),Inches(4),Inches(0.4))
    tf = txb.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = '--- Confidencial ---'
    r.font.size=Pt(10); r.font.color.rgb=BLANCO; r.font.italic=True

def add_linea_h(slide, y=Inches(2.1)):
    line = slide.shapes.add_shape(1,Inches(0.5),y,Inches(12.3),Inches(0.02))
    line.fill.solid(); line.fill.fore_color.rgb=RGBColor(0xCC,0xCC,0xCC)
    line.line.fill.background()

def add_text(slide,text,x,y,w,h,size=18,bold=False,color=GRIS_TEXTO,align=PP_ALIGN.LEFT,italic=False):
    txb = slide.shapes.add_textbox(x,y,w,h)
    tf = txb.text_frame; tf.word_wrap=True
    p = tf.paragraphs[0]; p.alignment=align
    r = p.add_run(); r.text=text
    r.font.size=Pt(size); r.font.bold=bold
    r.font.color.rgb=color; r.font.italic=italic
    return txb

def generar_credito(d, tmpdir):
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H
    blank = prs.slide_layouts[6]

    exitoso = d.get('exitoso', True)

    # ── Slide 1: Portada ──
    s1 = prs.slides.add_slide(blank)
    fondo_rojo(s1)
    if os.path.exists(LOGO_TC):
        s1.shapes.add_picture(LOGO_TC, Inches(9.8), Inches(0.25), Inches(3.2), Inches(0.82))
    add_text(s1,'Tecnocontrol Vehicular',Inches(0.6),Inches(1.5),Inches(10),Inches(1.2),
             size=44,color=AZUL_TITULO)
    titulo_portada = f'Reporte de recuperación {"Exitosa" if exitoso else "No Exitosa"}\nTK: {d.get("tk","—")}'
    add_text(s1,titulo_portada,Inches(0.6),Inches(3.2),Inches(11),Inches(2),size=48,color=BLANCO)

    # ── Slide 2: Datos del caso ──
    s2 = prs.slides.add_slide(blank)
    fondo_blanco(s2); add_confidencial(s2); add_linea_h(s2,y=Inches(1.4))
    from datetime import datetime
    def calcular_dias(f_inicio, f_fin):
        try:
            d1 = datetime.strptime(f_inicio.strip(), '%d/%m/%Y')
            d2 = datetime.strptime(f_fin.strip(), '%d/%m/%Y')
            return str(abs((d2 - d1).days))
        except Exception:
            return '—'

    dias = calcular_dias(d.get('fecha_inicio',''), d.get('fecha_fin',''))
    unidad = f"{d.get('marca','—')} {d.get('modelo','—')}".strip('— ') or '—'

    campos = [
        ('Nombre de cliente:',                        d.get('nombre','—')),
        ('Unidad:',                                   unidad),
        ('Año:',                                      d.get('anio','—')),
        ('Placas:',                                   d.get('placas','—')),
        ('Estatus:',                                  'Recuperación Exitosa' if exitoso else 'No Exitosa'),
        ('Motivo de recuperación:',                   d.get('motivo','—')),
        ('Fecha de inicio de seguimiento:',           d.get('fecha_inicio','—')),
        ('Fecha de finalización de seguimiento:',     d.get('fecha_fin','—')),
        ('Localidad de gestión:',                     d.get('localidad','—')),
        ('Lugar de resguardo:',                       d.get('lugar_resguardo','—')),
        ('Días que duró proceso de recuperación:',    dias),
        ('Gestor asignado:',                          d.get('gestor','—')),
    ]
    y_start = Inches(1.55); row_h = Inches(0.38)
    for i,(label,valor) in enumerate(campos):
        y = y_start + i*row_h
        add_text(s2,label,Inches(0.5),y,Inches(4.0),row_h,size=13,bold=True,color=GRIS_TEXTO)
        add_text(s2,valor,Inches(4.6),y,Inches(8.5),row_h,size=13,color=GRIS_TEXTO)

    # ── Slide 3: Mapa ──
    s3 = prs.slides.add_slide(blank)
    fondo_blanco(s3); add_confidencial(s3)
    add_text(s3,'Resultado de rastreo de unidad en últimas ubicaciones:',
             Inches(0.5),Inches(0.3),Inches(12),Inches(0.8),size=24,color=GRIS_TEXTO)
    add_linea_h(s3,y=Inches(1.2))

    img_mapa_b64 = d.get('img_mapa_credito')
    if img_mapa_b64:
        import io as _io
        from PIL import Image as PILImage
        img_bytes = base64.b64decode(img_mapa_b64)
        img = PILImage.open(_io.BytesIO(img_bytes))
        img_tmp = os.path.join(tmpdir, '_mapa_credito.jpg')
        img.convert('RGB').save(img_tmp, 'JPEG', quality=92)
        # Escalar para que quepa en el área del slide
        iw, ih = img.size
        max_w, max_h = Inches(12.3), Inches(4.8)
        scale = min(max_w / iw, max_h / ih)
        dw, dh = iw * scale, ih * scale
        x = Inches(0.5) + (max_w - dw) / 2
        s3.shapes.add_picture(img_tmp, x, Inches(1.35), dw, dh)
    else:
        map_box = s3.shapes.add_shape(1,Inches(1.5),Inches(1.35),Inches(10),Inches(4.8))
        map_box.fill.solid(); map_box.fill.fore_color.rgb=RGBColor(0xEE,0xEE,0xEE)
        map_box.line.color.rgb=RGBColor(0xCC,0xCC,0xCC)
        add_text(s3,'[ Insertar captura del mapa GPS ]',Inches(1.5),Inches(3.5),Inches(10),Inches(0.8),
                 size=14,color=RGBColor(0xAA,0xAA,0xAA),align=PP_ALIGN.CENTER)

    ubicacion = d.get('ubicacion','—')
    coords = d.get('coordenadas','—')
    add_text(s3, f'Coordenadas: {coords}',
             Inches(0.5), Inches(6.0), Inches(12), Inches(0.4), size=13, color=GRIS_TEXTO)
    add_text(s3, f'Ubicación: {ubicacion}',
             Inches(0.5), Inches(6.4), Inches(12), Inches(0.4), size=13, color=GRIS_TEXTO)

    # ── Slides Resumen ──
    resumen = d.get('resumen', [])
    ROWS_PER_SLIDE = 5
    for chunk_i in range(0, max(len(resumen), 1), ROWS_PER_SLIDE):
        chunk = resumen[chunk_i:chunk_i+ROWS_PER_SLIDE]
        if not chunk:
            break
        s = prs.slides.add_slide(blank)
        fondo_blanco(s); add_confidencial(s)
        add_text(s,'Resumen:',Inches(0.5),Inches(0.2),Inches(5),Inches(0.5),size=20,color=GRIS_TEXTO)
        add_linea_h(s,y=Inches(0.85))
        hdr_y = Inches(0.92)
        hdr = s.shapes.add_shape(1,Inches(0.5),hdr_y,Inches(12.3),Inches(0.35))
        hdr.fill.solid(); hdr.fill.fore_color.rgb=RGBColor(0x1B,0x2A,0x4A)
        hdr.line.fill.background()
        add_text(s,'FECHA',Inches(0.55),hdr_y,Inches(2.5),Inches(0.35),size=11,bold=True,color=BLANCO)
        add_text(s,'ACCIÓN',Inches(3.1),hdr_y,Inches(9.7),Inches(0.35),size=11,bold=True,color=BLANCO)
        row_h = Inches(0.95)
        for i,(fecha,accion) in enumerate(chunk):
            y = hdr_y + Inches(0.35) + i*row_h
            fill_color = RGBColor(0xF5,0xF7,0xFA) if i%2==0 else RGBColor(0xFF,0xFF,0xFF)
            row_bg = s.shapes.add_shape(1,Inches(0.5),y,Inches(12.3),row_h)
            row_bg.fill.solid(); row_bg.fill.fore_color.rgb=fill_color
            row_bg.line.color.rgb=RGBColor(0xCC,0xCC,0xCC)
            add_text(s,fecha,Inches(0.6),y+Inches(0.05),Inches(2.4),row_h,size=10,bold=True,color=GRIS_TEXTO)
            add_text(s,accion,Inches(3.1),y+Inches(0.05),Inches(9.6),row_h,size=10,color=GRIS_TEXTO)

    # ── Slide Evidencia ──
    s_ev = prs.slides.add_slide(blank)
    fondo_blanco(s_ev); add_confidencial(s_ev)
    add_text(s_ev,'Evidencia de la unidad:',Inches(0.5),Inches(0.3),Inches(8),Inches(0.7),size=24,color=GRIS_TEXTO)
    add_linea_h(s_ev,y=Inches(1.1))

    imgs_ev = d.get('img_evidencias_credito', [])
    import io as _io
    from PIL import Image as PILImage

    for i in range(3):
        x = Inches(0.4 + i*4.3)
        if i < len(imgs_ev) and imgs_ev[i]:
            try:
                img_bytes = base64.b64decode(imgs_ev[i])
                img = PILImage.open(_io.BytesIO(img_bytes)).convert('RGB')
                iw, ih = img.size
                max_w, max_h = Inches(4.0), Inches(5.2)
                scale = min(max_w/iw, max_h/ih)
                dw, dh = iw*scale, ih*scale
                img_tmp = os.path.join(tmpdir, f'_ev_credito_{i}.jpg')
                img.save(img_tmp, 'JPEG', quality=92)
                s_ev.shapes.add_picture(img_tmp, x, Inches(1.25), dw, dh)
            except Exception:
                ph = s_ev.shapes.add_shape(1,x,Inches(1.25),Inches(4.0),Inches(5.2))
                ph.fill.solid(); ph.fill.fore_color.rgb=RGBColor(0xEE,0xEE,0xEE)
                ph.line.color.rgb=RGBColor(0xCC,0xCC,0xCC)
        else:
            ph = s_ev.shapes.add_shape(1,x,Inches(1.25),Inches(4.0),Inches(5.2))
            ph.fill.solid(); ph.fill.fore_color.rgb=RGBColor(0xEE,0xEE,0xEE)
            ph.line.color.rgb=RGBColor(0xCC,0xCC,0xCC)
            add_text(s_ev,f'[ Foto {i+1} ]',x,Inches(3.5),Inches(4.0),Inches(0.5),
                     size=12,color=RGBColor(0xAA,0xAA,0xAA),align=PP_ALIGN.CENTER)

    # ── Slide Conclusiones ──
    s_conc = prs.slides.add_slide(blank)
    fondo_blanco(s_conc); add_confidencial(s_conc)
    add_text(s_conc,'Conclusiones:',Inches(0.5),Inches(0.3),Inches(8),Inches(0.7),size=24,color=GRIS_TEXTO)
    add_linea_h(s_conc,y=Inches(1.1))
    conclusion = d.get('conclusion',
        'Tras el proceso de notificación y el período de gracia otorgado al propietario, se procedió con '
        'la ejecución de la recuperación siguiendo todos los procedimientos pertinentes. El vehículo ha '
        'sido asegurado y se encuentra en resguardo para su correcta custodia hasta la resolución final del asunto.\n\n'
        'La recuperación del vehículo no pagado se ha llevado a cabo de manera efectiva y conforme a los '
        'procesos establecidos.' if exitoso else
        'Se concluye que el caso deberá ser atendido por el área legal correspondiente de NMDP, derivado '
        'de que la unidad se encuentra sin reportar desde hace un tiempo considerable y no ha sido posible '
        'visualizarla en las visitas técnicas realizadas.\n\n'
        'Por lo anterior, y en conjunto con el personal de NMDP, se determina dar por finalizado el protocolo '
        'de recuperación, a fin de que el área correspondiente proceda conforme a sus atribuciones.'
    )
    add_text(s_conc, conclusion, Inches(0.5), Inches(1.3), Inches(12.3), Inches(4), size=16, color=GRIS_TEXTO)

    # ── Guardar y convertir ──
    folio = d.get('folio','REPORTE')
    pptx_path = os.path.join(tmpdir, f'{folio}.pptx')
    pdf_path  = os.path.join(tmpdir, f'{folio}.pdf')
    prs.save(pptx_path)

    subprocess.run(
        ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', tmpdir, pptx_path],
        check=True,
        env={**os.environ, 'HOME': tmpdir}
    )
    return pdf_path
